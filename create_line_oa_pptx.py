from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Bosch brand colors
BOSCH_RED    = RGBColor(0xE2, 0x00, 0x15)
BOSCH_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
BOSCH_DGRAY  = RGBColor(0x4A, 0x4A, 0x4A)
BOSCH_MGRAY  = RGBColor(0x76, 0x76, 0x76)
BOSCH_LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
BOSCH_BGRAY  = RGBColor(0xEB, 0xEB, 0xEB)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def tb(slide, text, l, t, w, h, size, bold=False, color=BOSCH_BLACK,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


# ── shared chrome ──────────────────────────────────────────────────────────────
def content_chrome(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)

    # Left red sidebar
    rect(slide, 0, 0, 0.18, 7.5, BOSCH_RED)

    # Top black bar
    rect(slide, 0.18, 0, 13.15, 1.15, BOSCH_BLACK)

    # Bosch logo "circle" (red circle on black bar)
    rect(slide, 11.8, 0.18, 0.78, 0.78, BOSCH_RED)
    tb(slide, "BOSCH", 11.75, 0.22, 0.9, 0.6, 7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide title in top bar
    tb(slide, title, 0.5, 0.2, 11.0, 0.75, 24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Thin red rule under the bar
    rect(slide, 0.18, 1.15, 13.15, 0.06, BOSCH_RED)

    # Bottom footer bar
    rect(slide, 0.18, 7.1, 13.15, 0.4, BOSCH_LGRAY)
    tb(slide, "WhiteWings  |  LINE OA Solutions  |  2026",
       0.4, 7.12, 12.7, 0.3, 9, bold=False, color=BOSCH_MGRAY, align=PP_ALIGN.LEFT)

    return slide


def bullet_slide(prs, title, bullets):
    slide = content_chrome(prs, title)
    y = 1.42
    for i, (btitle, bdesc) in enumerate(bullets):
        # Red numbered square
        rect(slide, 0.35, y, 0.42, 0.42, BOSCH_RED)
        tb(slide, str(i + 1), 0.35, y + 0.02, 0.42, 0.38,
           13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Bullet title (black bold)
        tb(slide, btitle, 0.9, y + 0.01, 12.1, 0.38,
           13, bold=True, color=BOSCH_BLACK)
        y += 0.42

        # Bullet description (dark gray)
        if bdesc:
            tb(slide, bdesc, 0.9, y - 0.04, 12.1, 0.42,
               10.5, bold=False, color=BOSCH_DGRAY)
            y += 0.44

        # Thin separator line
        rect(slide, 0.35, y + 0.06, 12.65, 0.03, BOSCH_BGRAY)
        y += 0.2

    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1 · Cover ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, BOSCH_BLACK)

    # Full-height red left block
    rect(slide, 0, 0, 4.8, 7.5, BOSCH_RED)

    # White right panel
    rect(slide, 4.8, 0, 8.53, 7.5, WHITE)

    # Red accent strip at right-panel top
    rect(slide, 4.8, 0, 8.53, 0.12, BOSCH_RED)

    # Bosch badge on left panel
    rect(slide, 0.45, 0.45, 1.4, 1.4, WHITE)
    tb(slide, "BOSCH", 0.45, 0.62, 1.4, 0.8, 22, bold=True, color=BOSCH_RED, align=PP_ALIGN.CENTER)

    # Left panel — vertical label
    tb(slide, "TECHNOLOGY\nFOR BUSINESS", 0.3, 2.2, 4.0, 1.2,
       13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Horizontal rule on left panel
    rect(slide, 0.3, 3.55, 3.8, 0.06, WHITE)

    tb(slide, "WhiteWings", 0.3, 3.75, 4.0, 0.5,
       11, bold=False, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.LEFT)

    # Right panel — main title
    tb(slide, "LINE Official\nAccount", 5.2, 1.1, 7.7, 2.2,
       40, bold=True, color=BOSCH_BLACK, align=PP_ALIGN.LEFT)

    # Red underline beneath title
    rect(slide, 5.2, 3.35, 5.5, 0.1, BOSCH_RED)

    tb(slide, "(LINE OA) Benefits", 5.2, 3.6, 7.7, 0.7,
       22, bold=False, color=BOSCH_DGRAY, align=PP_ALIGN.LEFT)

    tb(slide, "Unlock the Power of Business Communication on LINE",
       5.2, 4.5, 7.7, 0.8, 13, bold=False, color=BOSCH_MGRAY, align=PP_ALIGN.LEFT)

    # Bottom footer strip on right panel
    rect(slide, 4.8, 6.9, 8.53, 0.6, BOSCH_LGRAY)
    tb(slide, "2026  |  Confidential", 5.0, 7.0, 8.0, 0.4,
       10, bold=False, color=BOSCH_MGRAY, align=PP_ALIGN.LEFT)

    # ── Slide 2 · What is LINE OA? ───────────────────────────────────────────
    slide2 = content_chrome(prs, "What is LINE Official Account?")

    tb(slide2,
       "LINE Official Account (LINE OA) is a business account on the LINE messaging platform "
       "that enables companies, brands, and organizations to communicate directly with their "
       "customers at scale. With over 90 million active users in Thailand alone, LINE OA "
       "provides an unparalleled channel to reach, engage, and convert your audience.",
       0.4, 1.35, 12.5, 1.3, 13, bold=False, color=BOSCH_DGRAY)

    # Stat cards
    stats = [
        ("200M+", "Monthly Active\nUsers (Global)"),
        ("90M+",  "Users in\nThailand"),
        ("#1",    "Messaging App\nin Thailand"),
        ("94%",   "Market\nPenetration"),
    ]
    x = 0.35
    for val, label in stats:
        # Card background
        rect(slide2, x, 2.9, 2.95, 2.9, BOSCH_LGRAY)
        # Red top accent
        rect(slide2, x, 2.9, 2.95, 0.12, BOSCH_RED)
        tb(slide2, val,   x, 3.15, 2.95, 1.0, 32, bold=True,  color=BOSCH_RED,   align=PP_ALIGN.CENTER)
        tb(slide2, label, x, 4.2,  2.95, 0.9, 11, bold=False, color=BOSCH_DGRAY, align=PP_ALIGN.CENTER)
        x += 3.25

    tb(slide2, "Source: LINE Corporation, 2024",
       0.4, 6.05, 12.0, 0.35, 9, bold=False, color=BOSCH_MGRAY)

    # ── Slide 3 · Core Benefits ───────────────────────────────────────────────
    bullet_slide(prs, "Core Benefits of LINE OA", [
        ("Direct Messaging to Customers",
         "Send personalized messages, promotions, and updates directly to followers' chat inbox — ensuring high open rates."),
        ("Broadcast & Segmented Messaging",
         "Reach all followers at once or target specific audience segments based on demographics, behavior, or tags."),
        ("Rich Menu & Custom UI",
         "Design interactive menus at the bottom of chats to guide users to key features, websites, or services instantly."),
        ("Chatbot & Auto-Reply Integration",
         "Automate responses with AI chatbots, FAQs, and triggered messages to provide 24/7 customer support."),
        ("Verified Badge & Brand Trust",
         "Earn the official verified badge that builds brand credibility and distinguishes you from personal accounts."),
    ])

    # ── Slide 4 · Marketing & Sales ──────────────────────────────────────────
    bullet_slide(prs, "Marketing & Sales Benefits", [
        ("Coupon & Reward Distribution",
         "Issue digital coupons, stamp cards, and loyalty rewards directly through LINE to drive repeat purchases."),
        ("LINE Points & Incentives",
         "Reward customers with LINE Points to encourage engagement, referrals, and brand loyalty."),
        ("Product Catalog & In-Chat Shopping",
         "Showcase products with images, descriptions, and links — enabling seamless in-app discovery and purchase."),
        ("Click-to-LINE Ads Integration",
         "Run targeted LINE Ads that direct users straight to a conversation with your OA for instant lead capture."),
        ("Analytics & Insights Dashboard",
         "Track message open rates, follower growth, chat interactions, and campaign performance in real time."),
    ])

    # ── Slide 5 · Customer Service ────────────────────────────────────────────
    bullet_slide(prs, "Customer Service Benefits", [
        ("Multi-Agent Chat Management",
         "Multiple staff members can manage and respond to customer chats simultaneously through one OA inbox."),
        ("Chat Tags & Customer Labels",
         "Organize conversations with custom tags and labels for efficient customer tracking and follow-up."),
        ("Automated Welcome Messages",
         "Greet every new follower instantly with a personalized welcome message and onboarding flow."),
        ("Appointment & Booking Systems",
         "Integrate booking modules so customers can schedule appointments directly within the LINE conversation."),
        ("Order Status & Notifications",
         "Send transactional notifications like order confirmations, shipping updates, and payment receipts via chat."),
    ])

    # ── Slide 6 · Technical Integration ──────────────────────────────────────
    bullet_slide(prs, "Technical & Integration Benefits", [
        ("LINE Login & Social Authentication",
         "Enable users to log in to your website or app using their LINE account for a frictionless experience."),
        ("Messaging API & Webhooks",
         "Use the LINE Messaging API to build custom bots, integrate CRM systems, and automate workflows."),
        ("LIFF (LINE Front-end Framework)",
         "Build mini web apps that run inside LINE, delivering rich interactive experiences without leaving the app."),
        ("CRM & Third-Party Integration",
         "Connect LINE OA with Salesforce, HubSpot, Zapier, and other tools for unified customer data management."),
        ("LINE Pay Integration",
         "Accept payments seamlessly within LINE conversations using LINE Pay for a smooth checkout experience."),
    ])

    # ── Slide 7 · Plans & Pricing ─────────────────────────────────────────────
    slide7 = content_chrome(prs, "Plans & Pricing Overview")

    plans = [
        ("Free Plan",     "0 THB / mo",
         ["500 messages/month", "Basic chat features", "1 linked account", "Standard chat support"]),
        ("Light Plan",    "~1,200 THB / mo",
         ["15,000 messages/month", "Rich Menu", "Multi-admin access", "Analytics dashboard"]),
        ("Standard Plan", "~3,900 THB / mo",
         ["45,000 messages/month", "All advanced features", "Full API access", "Priority support"]),
    ]
    x = 0.35
    for name, price, feats in plans:
        # Card
        rect(slide7, x, 1.38, 4.1, 5.5, BOSCH_LGRAY)
        # Red header strip
        rect(slide7, x, 1.38, 4.1, 1.1, BOSCH_RED)
        tb(slide7, name,  x, 1.42, 4.1, 0.55, 15, bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
        tb(slide7, price, x, 1.93, 4.1, 0.48, 13, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

        fy = 2.65
        for feat in feats:
            # Small red tick square
            rect(slide7, x + 0.18, fy + 0.05, 0.15, 0.28, BOSCH_RED)
            tb(slide7, feat, x + 0.45, fy, 3.5, 0.42, 11.5, bold=False, color=BOSCH_DGRAY)
            fy += 0.52

        x += 4.45

    tb(slide7, "* Pricing may vary. Please visit line.biz for the latest rates.",
       0.4, 7.0, 12.5, 0.35, 9, bold=False, color=BOSCH_MGRAY)

    # ── Slide 8 · Summary / CTA ───────────────────────────────────────────────
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide8, WHITE)

    # Red left sidebar
    rect(slide8, 0, 0, 0.18, 7.5, BOSCH_RED)
    # Black top bar (full width)
    rect(slide8, 0.18, 0, 13.15, 1.1, BOSCH_BLACK)
    # Red accent rule
    rect(slide8, 0.18, 1.1, 13.15, 0.06, BOSCH_RED)

    # Bosch badge
    rect(slide8, 11.8, 0.16, 0.78, 0.78, BOSCH_RED)
    tb(slide8, "BOSCH", 11.75, 0.2, 0.9, 0.6, 7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tb(slide8, "Why Choose LINE OA?",
       0.5, 0.18, 11.0, 0.75, 24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    points = [
        "Reach Thailand's #1 messaging platform with 90M+ users",
        "Drive sales with coupons, rewards, and in-chat shopping",
        "Automate customer support with chatbots and 24/7 auto-replies",
        "Gain real-time insights with powerful analytics and reporting tools",
        "Integrate seamlessly with your existing CRM and business systems",
    ]
    y = 1.35
    for pt in points:
        rect(slide8, 0.35, y + 0.06, 0.38, 0.38, BOSCH_RED)
        tb(slide8, str(points.index(pt) + 1), 0.35, y + 0.07, 0.38, 0.35,
           12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(slide8, pt, 0.9, y + 0.05, 12.0, 0.45,
           14, bold=False, color=BOSCH_BLACK)
        rect(slide8, 0.35, y + 0.52, 12.65, 0.03, BOSCH_BGRAY)
        y += 0.75

    # CTA box
    rect(slide8, 0.35, 5.55, 12.65, 1.1, BOSCH_RED)
    tb(slide8, "Get Started Today", 0.5, 5.62, 12.33, 0.55,
       22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide8, "line.biz  |  Contact WhiteWings for LINE OA setup and integration",
       0.5, 6.1, 12.33, 0.45, 12, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Footer
    rect(slide8, 0.18, 7.1, 13.15, 0.4, BOSCH_LGRAY)
    tb(slide8, "WhiteWings  |  LINE OA Solutions  |  2026",
       0.4, 7.12, 12.7, 0.3, 9, bold=False, color=BOSCH_MGRAY, align=PP_ALIGN.LEFT)

    # Save
    out = "/home/user/WhiteWings/LINE_OA_Benefits.pptx"
    prs.save(out)
    print(f"Saved: {out}")


create_presentation()
